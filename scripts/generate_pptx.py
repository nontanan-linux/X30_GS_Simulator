import json
import os
import glob
import ast
import csv
from collections import Counter
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from docx import Document
from docx.shared import Inches as DocxInches, Pt as DocxPt, RGBColor as DocxRGBColor
from docx.oxml import OxmlElement
from docx.oxml.ns import qn

# Layout Constants
SLIDE_WIDTH = 13.333
SLIDE_HEIGHT = 7.5

# Colors
BG_COLOR = RGBColor(10, 15, 30)
CARD_BG = RGBColor(25, 35, 55)
ACCENT_BLUE = RGBColor(99, 102, 241)
TEXT_PRIMARY = RGBColor(255, 255, 255)
TEXT_SECONDARY = RGBColor(160, 174, 192)
# Project paths
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
PROJECT_ROOT = os.path.dirname(SCRIPT_DIR)

MISSION_ID = 'mission-69'
MISSION_DIR = os.path.join(PROJECT_ROOT, 'resource', MISSION_ID)
LOGO_PATH = os.path.join(PROJECT_ROOT, 'resource', 'gensurv-logo.jpg')

def find_images(point_name, inspection_type):
    suffixes = []
    if inspection_type == 'thermal_inspection':
        suffixes = ['_raw_rgb_image.jpg', '_thermal_map.jpg']
    elif inspection_type == 'gauge_inspection':
        suffixes = ['_preprocessed_rgb_image.jpg', '_processed_rgb_image.jpg']
    elif inspection_type in ['leakage_inspection', 'vibration_inspection']:
        suffixes = ['_spect_photo.jpg', '_voiceprint_photo.jpg', '_soundmap_photo.jpg']
    elif inspection_type in ['asset_inspection', 'loto_inspection']:
        suffixes = ['_processed_rgb_image.jpg']
    
    found_files = []
    for suffix in suffixes:
        pattern = os.path.join(MISSION_DIR, f'*{point_name}*{suffix}')
        matches = glob.glob(pattern)
        if matches:
            found_files.append(matches[0]) # Take first match
            
    return found_files

def add_slide(prs, point, index, total):
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR
    
    def add_text(left, top, width, height, text, font_size=12, bold=False, color=TEXT_PRIMARY, alignment=PP_ALIGN.LEFT):
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = alignment
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = 'Arial'
        return txBox

    def add_card(left, top, width, height, label, value):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = CARD_BG
        shape.line.visible = False
        add_text(left + 0.1, top + 0.1, width - 0.2, 0.3, label, font_size=9, color=TEXT_SECONDARY)
        add_text(left + 0.1, top + 0.35, width - 0.2, 0.4, value, font_size=14, bold=True)

    name = point.get('Node_info', 'N/A')
    raw_type = point.get('Inspection', 'General')
    type_display = raw_type.replace('_', ' ').title()
    zone = point.get('Zone', 'N/A')
    floor = point.get('MapName', 'N/A').replace('_', ' ').title()

    # 1. Main Title
    seq_num = f"{index+1:03d}"
    title_text = f"Inspection Point Details - [{seq_num}]"
    add_text(0.5, 0.4, 11.0, 0.8, title_text, font_size=36, bold=True)
    
    # 2. Logo (Flush Top-Right)
    logo_path = LOGO_PATH
    if os.path.exists(logo_path):
        logo_h = 0.6
        slide.shapes.add_picture(logo_path, Inches(12.733), Inches(0.0), height=Inches(logo_h))
    
    # 3. Metadata Row (Compacted)
    card_y = 1.3
    card_h = 1.0
    spacing = 0.1
    
    # New Widths
    w_name = 2.4
    w_type = 2.4
    w_loc = 2.0  # Combined Zone & Floor
    
    add_card(0.5, card_y, w_name, card_h, "Inspection Point Name", name)
    add_card(0.5 + w_name + spacing, card_y, w_type, card_h, "Inspection Type", type_display)
    
    # Combined Zone & Floor Card
    loc_x = 0.5 + w_name + w_type + spacing*2
    shape_loc = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(loc_x), Inches(card_y), Inches(w_loc), Inches(card_h))
    shape_loc.fill.solid()
    shape_loc.fill.fore_color.rgb = CARD_BG
    shape_loc.line.visible = False
    add_text(loc_x + 0.1, card_y + 0.05, w_loc - 0.2, 0.3, "Zone & Floor", font_size=9, color=TEXT_SECONDARY)
    add_text(loc_x + 0.1, card_y + 0.3, w_loc - 0.2, 0.3, f"Zone: {zone}", font_size=11, bold=True)
    add_text(loc_x + 0.1, card_y + 0.6, w_loc - 0.2, 0.3, f"Floor: {floor}", font_size=11, bold=True)
    
    # Expanded Result Status Card
    status_x = loc_x + w_loc + spacing
    w_status = (SLIDE_WIDTH - 0.5) - status_x
    
    threshold = point.get('Threshold', 'N/A')
    
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(status_x), Inches(card_y), Inches(w_status), Inches(card_h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.visible = False
    add_text(status_x + 0.1, card_y + 0.05, w_status - 0.2, 0.3, "Result Status", font_size=9, color=TEXT_SECONDARY)
    
    # Get Data from CSV if available
    notif = notifications.get(name, {})
    csv_res = notif.get('result_dict', {})
    csv_level = notif.get('notification_level', 'N/A').title()
    
    if raw_type == 'thermal_inspection':
        min_temp = csv_res.get('min_temperature', 'N/A')
        max_temp = csv_res.get('max_temperature', 'N/A')
        add_text(status_x + 0.1, card_y + 0.25, w_status - 0.2, 0.25, f"Notification Level: {csv_level}", font_size=10, bold=True)
        add_text(status_x + 0.1, card_y + 0.50, w_status - 0.2, 0.25, f"Threshold: {threshold}", font_size=10, bold=True)
        add_text(status_x + 0.1, card_y + 0.75, w_status - 0.2, 0.25, f"Temperature-Range: {min_temp}-{max_temp}", font_size=10, bold=True)
    elif raw_type == 'leakage_inspection':
        is_leak = csv_res.get('is_leak', 'N/A')
        add_text(status_x + 0.1, card_y + 0.25, w_status - 0.2, 0.25, f"Is Leak: {is_leak}", font_size=10, bold=True)
        add_text(status_x + 0.1, card_y + 0.50, w_status - 0.2, 0.25, f"Notification Level: {csv_level}", font_size=10, bold=True)
        add_text(status_x + 0.1, card_y + 0.75, w_status - 0.2, 0.25, "Frequency Range:", font_size=10, bold=True)
    elif raw_type == 'vibration_inspection':
        k_thresh = point.get('Kurtosis_threshold', 'N/A')
        p_thresh = point.get('Peak_threshold', 'N/A')
        k_val = csv_res.get('kurtosis', 'N/A')
        p_val = csv_res.get('peak_factor', 'N/A')
        add_text(status_x + 0.1, card_y + 0.25, w_status - 0.2, 0.2, f"Notification Level: {csv_level}", font_size=9, bold=True)
        add_text(status_x + 0.1, card_y + 0.43, w_status - 0.2, 0.2, "Frequency Range:", font_size=9, bold=True)
        add_text(status_x + 0.1, card_y + 0.61, w_status - 0.2, 0.2, f"Kurtosis Threshold: {k_thresh}, Kurtosis Detection: {k_val}", font_size=9, bold=True)
        add_text(status_x + 0.1, card_y + 0.79, w_status - 0.2, 0.2, f"Peak Factor Threshold: {p_thresh}, Peak Factor: {p_val}", font_size=9, bold=True)
    else:
        det_status = csv_res.get('detection_status', 'N/A')
        add_text(status_x + 0.1, card_y + 0.3, w_status - 0.2, 0.3, f"Detection Status: {det_status}", font_size=11, bold=True)
        add_text(status_x + 0.1, card_y + 0.6, w_status - 0.2, 0.3, f"Notification Level: {csv_level}", font_size=11, bold=True)

    # 4. Image Frames
    frame_y = 2.5
    
    # Custom widths for Leakage/Vibration (3 images)
    if raw_type in ['leakage_inspection', 'vibration_inspection']:
        f1_w = 4.8
        f2_w = 7.5
    else:
        f1_w = 6.0
        f2_w = 6.0
    
    frame_h = 3.5
    
    # Left Frame (Robot/Map Position Image from extracted PDF)
    f1_rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(frame_y), Inches(f1_w), Inches(frame_h))
    f1_rect.fill.solid()
    f1_rect.fill.fore_color.rgb = CARD_BG
    f1_rect.line.visible = False
    
    takescreen_dir = os.path.join(PROJECT_ROOT, 'TakeScreen')
    robot_img_path = os.path.join(takescreen_dir, f'{name}.jpg')
    
    if os.path.exists(robot_img_path):
        try:
            slide.shapes.add_picture(robot_img_path, Inches(0.5), Inches(frame_y), width=Inches(f1_w), height=Inches(frame_h))
        except Exception as e:
            print(f"Error adding robot image {robot_img_path}: {e}")
            add_text(0.5, frame_y + (frame_h/2) - 0.2, f1_w, 0.4, "Robot Image Corrupt", font_size=10, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)
    else:
        add_text(0.5, frame_y + (frame_h/2) - 0.2, f1_w, 0.4, "Robot & Map Image Not Found", font_size=10, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)

    # Right Frame (Inspection Results)
    f2_x = SLIDE_WIDTH - 0.5 - f2_w
    f2_rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(f2_x), Inches(frame_y), Inches(f2_w), Inches(frame_h))
    f2_rect.fill.solid()
    f2_rect.fill.fore_color.rgb = CARD_BG
    f2_rect.line.visible = False
    
    images = find_images(name, raw_type)
    if not images:
        add_text(f2_x, frame_y + (frame_h/2) - 0.2, f2_w, 0.4, "Inspection Result Image Not Found", font_size=12, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)
    else:
        num_imgs = len(images)
        img_w = f2_w / num_imgs if num_imgs > 0 else f2_w
        for i, img_path in enumerate(images):
            left_offset = f2_x + (i * img_w)
            try:
                slide.shapes.add_picture(img_path, Inches(left_offset), Inches(frame_y), width=Inches(img_w), height=Inches(frame_h))
            except Exception as e:
                print(f"Error adding image {img_path}: {e}")

    # 5. Description Area (Now with Dynamic Discussion)
    desc_y = 6.2
    add_text(0.5, desc_y, 2.0, 0.3, "Description / Discussion", font_size=10, color=TEXT_SECONDARY)
    
    discussion_text = get_discussion(name, raw_type, csv_res, csv_level)
    add_text(0.5, desc_y + 0.25, 12.0, 0.8, discussion_text, font_size=11, color=TEXT_PRIMARY)

    # 6. Pagination Dots
    dot_y = 7.1
    dot_size = 0.06
    for i in range(3):
        dot_color = ACCENT_BLUE if i == 0 else RGBColor(50, 60, 80)
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5 + (i * 0.15)), Inches(dot_y), Inches(dot_size), Inches(dot_size))
        dot.fill.solid()
        dot.fill.fore_color.rgb = dot_color
        dot.line.visible = False

# Load Notification Data from CSV
notifications = {}

def load_notifications():
    global notifications
    notifications.clear()
    csv_pattern = os.path.join(MISSION_DIR, 'notification_*.csv')
    csv_files = glob.glob(csv_pattern)
    csv_path = csv_files[0] if csv_files else os.path.join(PROJECT_ROOT, 'resource', 'mission-332', 'notification_332_2026-03-18_2026-04-24.csv')
    if os.path.exists(csv_path):
        with open(csv_path, mode='r', encoding='utf-8') as f:
            reader = csv.DictReader(f)
            for row in reader:
                name = row['action_name']
                # Parse result string to dict safely
                try:
                    res_str = row['result'].replace("'", '"').replace("True", "true").replace("False", "false")
                    row['result_dict'] = json.loads(res_str)
                except:
                    try:
                        row['result_dict'] = ast.literal_eval(row['result'])
                    except:
                        row['result_dict'] = {}
                notifications[name] = row

def get_discussion(name, raw_type, csv_res, csv_level):
    level = csv_level.lower()
    
    if raw_type == 'thermal_inspection':
        max_t = csv_res.get('max_temperature', 0)
        # We don't have the threshold in this function but we can use generic
        if level == 'pass':
            return "The thermal camera was utilized to monitor the surface temperature of the equipment. The test results indicate that the operating temperature is within the normal range."
        else:
            return f"The thermal camera detected an elevated temperature of {max_t}°C, which may indicate overheating or abnormal operating conditions."
            
    elif raw_type == 'leakage_inspection':
        is_leak = csv_res.get('is_leak', False)
        if is_leak:
            return "The Acoustic Imager successfully detected and localized an air leak, pinpointing its exact coordinates as shown in the visual overlay."
        else:
            return "The Acoustic Imager monitored the area for ultrasonic leak patterns. No significant air or gas leaks were detected, indicating the system is airtight."
            
    elif 'gauge' in raw_type:
        status = csv_res.get('detection_status', True)
        if not status:
            return "Due to the absence of clear indicator markings or color-coded safe ranges on the gauge, the system could not verify the status, resulting in a critical detection output."
        else:
            return "The vision system successfully identified the gauge reading. The indicator is positioned within the standard operating range."
            
    elif 'vibration' in raw_type:
        if level == 'pass':
            return "The system analyzed acoustic vibration patterns to assess the equipment's operational condition. The results indicate that it is functioning normally with no signs of mechanical damage."
        else:
            return "Acoustic vibration analysis detected irregular frequency patterns or high peak factors, suggesting potential mechanical wear or bearing issues."
            
    elif 'asset' in raw_type or 'loto' in raw_type:
        status = csv_res.get('detection_status', True)
        if not status:
            return "The vision system identified an anomaly in the designated area, detecting that certain equipment or protective assets were missing or not in their correct state."
        else:
            return "The vision system confirmed that all assets and safety indicators (LOTO) are in their correct positions and states."
            
    return f"The automated inspection for {name} was completed. The system evaluated the {raw_type} data and the results were recorded as {level}."

def add_issues_slide(prs, failed_points, missing_points):
    if not failed_points and not missing_points:
        return
        
    def create_slide():
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR
        
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12), Inches(0.8))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = "Action Required: Notifications & Missing Data"
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = TEXT_PRIMARY
        
        logo_path = LOGO_PATH
        if os.path.exists(logo_path):
            slide.shapes.add_picture(logo_path, Inches(12.733), Inches(0.0), height=Inches(0.6))
            
        return slide

    def add_text_block(slide, left, top, width, height, title, items, title_color):
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = title_color
        
        for item in items:
            p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(12)
            p.font.color.rgb = TEXT_PRIMARY
            
    max_items_per_col = 15
    cols = []
    
    failed_chunks = [failed_points[i:i + max_items_per_col] for i in range(0, len(failed_points), max_items_per_col)]
    for idx, chunk in enumerate(failed_chunks):
        title = f"NOTIFICATIONS ({len(failed_points)})" if idx == 0 else "NOTIFICATIONS (Cont.)"
        cols.append((title, chunk, RGBColor(248, 113, 113)))
        
    missing_chunks = [missing_points[i:i + max_items_per_col] for i in range(0, len(missing_points), max_items_per_col)]
    for idx, chunk in enumerate(missing_chunks):
        title = f"MISSING DATA ({len(missing_points)})" if idx == 0 else "MISSING DATA (Cont.)"
        cols.append((title, chunk, RGBColor(250, 204, 21)))
        
    for i in range(0, len(cols), 2):
        slide = create_slide()
        col1 = cols[i]
        add_text_block(slide, 1.0, 1.5, 5.5, 5.5, col1[0], col1[1], col1[2])
        
        if i + 1 < len(cols):
            col2 = cols[i+1]
            add_text_block(slide, 7.0, 1.5, 5.5, 5.5, col2[0], col2[1], col2[2])

def add_confusion_matrix_slide(prs, inspection_points, notifications):
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR
    
    def add_text(left, top, width, height, text, font_size=12, bold=False, color=TEXT_PRIMARY, alignment=PP_ALIGN.LEFT, font_name='Arial'):
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        lines = text.split('\n')
        for idx, line in enumerate(lines):
            if idx == 0:
                p = tf.paragraphs[0]
                p.text = line
            else:
                p = tf.add_paragraph()
                p.text = line
            p.alignment = alignment
            p.font.size = Pt(font_size)
            p.font.bold = bold
            p.font.color.rgb = color
            p.font.name = font_name
        return txBox

    # Title
    add_text(0.5, 0.4, 12, 0.8, "Confusion Matrix & Performance Evaluation", font_size=36, bold=True, alignment=PP_ALIGN.CENTER)
    
    # Logo
    logo_path = LOGO_PATH
    if os.path.exists(logo_path):
        slide.shapes.add_picture(logo_path, Inches(12.733), Inches(0.0), height=Inches(0.6))

    # Calculate Confusion Matrix
    tp = fp = fn = tn = 0
    for point in inspection_points:
        name = point.get('Node_info', '')
        raw_insp = point.get('Inspection', '')
        
        notif = notifications.get(name)
        if notif:
            level = notif.get('notification_level', 'pass').lower()
            csv_res = notif.get('result_dict', {})
        else:
            level = 'pass'
            csv_res = {}
            
        is_critical = (level == 'critical')
        
        is_thermal_fn = False
        if raw_insp == 'thermal_inspection':
            threshold = point.get('Threshold')
            if threshold is not None:
                try:
                    threshold_val = float(threshold)
                except:
                    threshold_val = 50.0
            else:
                threshold_val = 50.0
                
            max_temp = csv_res.get('max_temperature')
            if max_temp is not None:
                try:
                    max_temp_val = float(max_temp)
                except:
                    max_temp_val = 0.0
            else:
                max_temp_val = 0.0
                
            if max_temp_val > threshold_val:
                if not is_critical:
                    is_thermal_fn = True
                    
        if is_critical:
            tp += 1
        elif is_thermal_fn:
            fn += 1
        else:
            tn += 1

    # Calculate Metrics
    total = tp + tn + fp + fn
    accuracy = (tp + tn) / total if total > 0 else 0.0
    precision = tp / (tp + fp) if (tp + fp) > 0 else 0.0
    recall = tp / (tp + fn) if (tp + fn) > 0 else 0.0
    f1 = 2 * precision * recall / (precision + recall) if (precision + recall) > 0 else 0.0

    # Draw Confusion Matrix Table
    rows = 3
    cols = 3
    left_margin = Inches(1.0)
    top_margin = Inches(1.3)
    table_w = Inches(5.3)
    table_h = Inches(2.4)
    
    table_shape = slide.shapes.add_table(rows, cols, left_margin, top_margin, table_w, table_h)
    table = table_shape.table
    
    # Column widths
    table.columns[0].width = Inches(1.9)
    table.columns[1].width = Inches(1.7)
    table.columns[2].width = Inches(1.7)
    
    # Table data
    table_data = [
        ["Predicted \\ Actual", "Actual Normal", "Actual Anomaly"],
        ["Predicted Normal", f"TN\n{tn}", f"FN\n{fn}"],
        ["Predicted Anomaly", f"FP\n{fp}", f"TP\n{tp}"]
    ]
    
    # Colors
    color_header_bg = RGBColor(40, 50, 70)
    color_tn_tp_bg = RGBColor(20, 83, 45)     # Dark Green
    color_fn_fp_bg = RGBColor(127, 29, 29)    # Dark Red
    
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.fill.solid()
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            # Style backgrounds
            if r == 0 or c == 0:
                cell.fill.fore_color.rgb = color_header_bg
            elif (r == 1 and c == 1) or (r == 2 and c == 2):
                cell.fill.fore_color.rgb = color_tn_tp_bg
            else:
                cell.fill.fore_color.rgb = color_fn_fp_bg
                
            cell.text = table_data[r][c]
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = TEXT_PRIMARY
            p.font.name = 'Arial'
            p.alignment = PP_ALIGN.CENTER
            
            # If r > 0 and c > 0, style the label and value
            if r > 0 and c > 0:
                cell.text = ""
                # Add Label paragraph
                p_lbl = cell.text_frame.paragraphs[0]
                p_lbl.text = "TN" if (r==1 and c==1) else "FN" if (r==1 and c==2) else "FP" if (r==2 and c==1) else "TP"
                p_lbl.font.size = Pt(10)
                p_lbl.font.bold = False
                p_lbl.font.color.rgb = TEXT_SECONDARY
                p_lbl.font.name = 'Arial'
                p_lbl.alignment = PP_ALIGN.CENTER
                
                # Add Value paragraph
                p_val = cell.text_frame.add_paragraph()
                p_val.text = str(tn if (r==1 and c==1) else fn if (r==1 and c==2) else fp if (r==2 and c==1) else tp)
                p_val.font.size = Pt(24)
                p_val.font.bold = True
                p_val.font.color.rgb = TEXT_PRIMARY
                p_val.font.name = 'Arial'
                p_val.alignment = PP_ALIGN.CENTER

    # TP/TN/FP/FN Legend Box below table
    legend_tp_tn = (
        "• TP (True Positive)  : Anomaly correctly detected\n"
        "• TN (True Negative)  : Normal condition correctly verified\n"
        "• FP (False Positive) : Normal wrongly flagged as anomaly\n"
        "• FN (False Negative) : Anomaly missed (predicted normal)"
    )
    add_text(1.0, 3.9, 5.3, 0.9, legend_tp_tn, font_size=8.5, color=TEXT_PRIMARY, font_name='Consolas')

    # Draw Metrics Cards on the right
    card_w = 2.4
    card_h = 0.85
    c1_x = 7.0
    c2_x = 9.9
    r1_y = 1.3
    r2_y = 2.35
    
    def add_metric_card(left, top, label, val_text):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(card_w), Inches(card_h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = CARD_BG
        shape.line.visible = False
        add_text(left + 0.1, top + 0.08, card_w - 0.2, 0.25, label, font_size=9, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)
        add_text(left + 0.1, top + 0.3, card_w - 0.2, 0.45, val_text, font_size=20, bold=True, color=TEXT_PRIMARY, alignment=PP_ALIGN.CENTER)

    add_metric_card(c1_x, r1_y, "ACCURACY", f"{accuracy:.1%}")
    add_metric_card(c2_x, r1_y, "PRECISION", f"{precision:.1%}")
    add_metric_card(c1_x, r2_y, "RECALL (SENSITIVITY)", f"{recall:.1%}")
    add_metric_card(c2_x, r2_y, "F1-SCORE", f"{f1:.1%}")

    # Metrics Definitions Box below the cards
    legend_metrics = (
        "• Accuracy  : Overall proportion of correct predictions\n"
        "• Precision : Alarm reliability (TP / Predicted Anomaly)\n"
        "• Recall    : Anomaly detection rate (TP / Actual Anomaly)\n"
        "• F1-Score  : Harmonic mean of Precision and Recall"
    )
    add_text(7.0, 3.9, 5.3, 0.9, legend_metrics, font_size=8.5, color=TEXT_PRIMARY, font_name='Consolas')

    # Description Area at the bottom (formatted equations)
    desc_y = 5.1
    add_text(1.0, desc_y, 11.3, 0.3, "Calculation Steps", font_size=11, color=TEXT_SECONDARY, bold=True)

    def format_fraction(label, num_str, den_str, val_num_str, val_den_str, result_str):
        max_top_len = max(len(num_str), len(val_num_str))
        max_bot_len = max(len(den_str), len(val_den_str))
        width = max(max_top_len, max_bot_len)
        
        num_padded = num_str.center(width)
        val_num_padded = val_num_str.center(width)
        
        den_padded = den_str.center(width)
        val_den_padded = val_den_str.center(width)
        
        bar = '-' * width
        
        lbl_len = len(label)
        lbl_spaces = ' ' * lbl_len
        
        line1 = f"{lbl_spaces}   {num_padded}     {val_num_padded}"
        line2 = f"{label} = {bar}  =  {bar}  =  {result_str}"
        line3 = f"{lbl_spaces}   {den_padded}     {val_den_padded}"
        
        return f"{line1}\n{line2}\n{line3}"

    eq_accuracy = format_fraction("• Accuracy ", "TP + TN", "Total", f"{tp} + {tn}", f"{total}", f"{accuracy:.1%}")
    eq_precision = format_fraction("• Precision", "TP", "TP + FP", f"{tp}", f"{tp} + {fp}", f"{precision:.1%}")
    eq_recall = format_fraction("• Recall   ", "TP", "TP + FN", f"{tp}", f"{tp} + {fn}", f"{recall:.1%}")
    eq_f1 = format_fraction("• F1-Score ", "2 * P * R", "P + R", f"2 * {precision:.2f} * {recall:.2f}", f"{precision:.2f} + {recall:.2f}", f"{f1:.1%}")

    col1_text = f"{eq_accuracy}\n\n{eq_precision}"
    col2_text = f"{eq_recall}\n\n{eq_f1}"

    # Draw formulas in two columns using monospace Consolas font
    add_text(1.0, desc_y + 0.4, 5.3, 1.6, col1_text, font_size=9.5, color=TEXT_PRIMARY, font_name='Consolas')
    add_text(7.0, desc_y + 0.4, 5.3, 1.6, col2_text, font_size=9.5, color=TEXT_PRIMARY, font_name='Consolas')

def add_summary_slide(prs, inspection_points, notifications):
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR
    
    def add_text(left, top, width, height, text, font_size=12, bold=False, color=TEXT_PRIMARY, alignment=PP_ALIGN.LEFT):
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = alignment
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = 'Arial'
        return txBox

    # Title
    add_text(0.5, 0.4, 12, 0.8, "Inspection Summary Report", font_size=40, bold=True, alignment=PP_ALIGN.CENTER)
    
    # Logo
    logo_path = LOGO_PATH
    if os.path.exists(logo_path):
        slide.shapes.add_picture(logo_path, Inches(12.733), Inches(0.0), height=Inches(0.6))

    # Aggregate Data
    summary_data = {} # {type: {level: count}}
    total_counts = {"pass": 0, "fail": 0, "missing": 0}
    missing_points = []
    failed_points = []
    
    for point in inspection_points:
        name = point.get('Node_info', '')
        raw_insp = point.get('Inspection', 'General')
        
        # Split LOTO from Asset if name contains 'loto'
        if raw_insp == 'asset_inspection' and 'loto' in name.lower():
            raw_type = 'Loto Inspection'
        else:
            raw_type = raw_insp.replace('_', ' ').title()
        
        notif = notifications.get(name)
        images = find_images(name, raw_insp)
        
        if not notif or not images:
            level = 'missing'
            if not notif and not images:
                missing_points.append(f"{name} (No CSV, No Images)")
            elif not notif:
                missing_points.append(f"{name} (No CSV)")
            elif not images:
                missing_points.append(f"{name} (No Images)")
        else:
            level = notif.get('notification_level', 'N/A').lower()
        
        if raw_type not in summary_data:
            summary_data[raw_type] = Counter()
        
        summary_data[raw_type][level] += 1
        
        if level == 'pass':
            total_counts['pass'] += 1
        elif level == 'missing':
            total_counts['missing'] += 1
        else:
            total_counts['fail'] += 1
            failed_points.append(f"{name} ({level.title()})")

    if missing_points:
        print(f"\n[Warning] Missing data detected for {len(missing_points)} points:")
        for mp in missing_points:
            print(f"  - {mp}")
        print()
        
    if failed_points:
        print(f"\n[Warning] Notifications detected for {len(failed_points)} points:")
        for fp in failed_points:
            print(f"  - {fp}")
        print()

    # Main Stats Cards
    card_y = 1.5
    card_w = 2.8
    card_h = 1.2
    spacing = 0.4
    
    start_x = (SLIDE_WIDTH - (4 * card_w + 3 * spacing)) / 2
    
    # Total Points Card
    t_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(start_x), Inches(card_y), Inches(card_w), Inches(card_h))
    t_card.fill.solid()
    t_card.fill.fore_color.rgb = CARD_BG
    t_card.line.visible = False
    add_text(start_x + 0.1, card_y + 0.1, card_w - 0.2, 0.3, "TOTAL POINTS", font_size=12, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)
    add_text(start_x + 0.1, card_y + 0.4, card_w - 0.2, 0.6, str(len(inspection_points)), font_size=36, bold=True, color=RGBColor(255, 255, 255), alignment=PP_ALIGN.CENTER)

    # Pass Card
    p_x = start_x + card_w + spacing
    p_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(p_x), Inches(card_y), Inches(card_w), Inches(card_h))
    p_card.fill.solid()
    p_card.fill.fore_color.rgb = CARD_BG
    p_card.line.visible = False
    add_text(p_x + 0.1, card_y + 0.1, card_w - 0.2, 0.3, "TOTAL PASS", font_size=12, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)
    add_text(p_x + 0.1, card_y + 0.4, card_w - 0.2, 0.6, str(total_counts['pass']), font_size=36, bold=True, color=RGBColor(74, 222, 128), alignment=PP_ALIGN.CENTER)

    # Fail/Notification Card
    f_x = p_x + card_w + spacing
    f_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(f_x), Inches(card_y), Inches(card_w), Inches(card_h))
    f_card.fill.solid()
    f_card.fill.fore_color.rgb = CARD_BG
    f_card.line.visible = False
    add_text(f_x + 0.1, card_y + 0.1, card_w - 0.2, 0.3, "NOTIFICATIONS", font_size=12, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)
    add_text(f_x + 0.1, card_y + 0.4, card_w - 0.2, 0.6, str(total_counts['fail']), font_size=36, bold=True, color=RGBColor(248, 113, 113), alignment=PP_ALIGN.CENTER)

    # Missing Card
    m_x = f_x + card_w + spacing
    m_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(m_x), Inches(card_y), Inches(card_w), Inches(card_h))
    m_card.fill.solid()
    m_card.fill.fore_color.rgb = CARD_BG
    m_card.line.visible = False
    add_text(m_x + 0.1, card_y + 0.1, card_w - 0.2, 0.3, "MISSING DATA", font_size=12, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)
    add_text(m_x + 0.1, card_y + 0.4, card_w - 0.2, 0.6, str(total_counts['missing']), font_size=36, bold=True, color=RGBColor(250, 204, 21), alignment=PP_ALIGN.CENTER)

    # Breakdown Table Header
    table_y = 3.2
    add_text(0.5, table_y, 4, 0.4, "Breakdown by Inspection Type", font_size=18, bold=True)
    
    col_w_type = 3.5
    col_w_pass = 1.8
    col_w_fail = 1.8
    col_w_miss = 1.8
    col_w_total = 1.8
    total_w = col_w_type + col_w_pass + col_w_fail + col_w_miss + col_w_total
    start_x = (SLIDE_WIDTH - total_w) / 2
    
    rows = len(summary_data) + 2  # Header row + data rows + Total row
    cols = 5
    left = Inches(start_x)
    top = Inches(table_y + 0.5)
    width = Inches(total_w)
    height = Inches(0.4 * rows)
    
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    
    # Set Column Widths
    table.columns[0].width = Inches(col_w_type)
    table.columns[1].width = Inches(col_w_pass)
    table.columns[2].width = Inches(col_w_fail)
    table.columns[3].width = Inches(col_w_miss)
    table.columns[4].width = Inches(col_w_total)
    
    # Draw Headers
    headers = ["Inspection Type", "Pass Count", "Notification", "Missing", "Total"]
    for c, text in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(40, 50, 70)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # Style text
        cell.text = text
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = TEXT_PRIMARY
        p.font.name = 'Arial'
        p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
        
    # Variables to hold column sums
    sum_pass = 0
    sum_fail = 0
    sum_miss = 0
    sum_total = 0

    # Draw Data Rows
    for r, (raw_type, counts) in enumerate(sorted(summary_data.items()), start=1):
        pass_count = counts.get('pass', 0)
        miss_count = counts.get('missing', 0)
        fail_count = sum(c for lvl, c in counts.items() if lvl not in ('pass', 'missing'))
        row_total = pass_count + fail_count + miss_count
        
        sum_pass += pass_count
        sum_fail += fail_count
        sum_miss += miss_count
        sum_total += row_total
        
        row_data = [raw_type, str(pass_count), str(fail_count), str(miss_count), str(row_total)]
        for c, text in enumerate(row_data):
            cell = table.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = CARD_BG
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            cell.text = text
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(14)
            p.font.bold = False
            p.font.color.rgb = TEXT_PRIMARY
            p.font.name = 'Arial'
            p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER

    # Draw Bottom Total Row
    total_row_idx = len(summary_data) + 1
    total_row_data = ["Total", str(sum_pass), str(sum_fail), str(sum_miss), str(sum_total)]
    logo_color = RGBColor(0, 38, 77)
    for c, text in enumerate(total_row_data):
        cell = table.cell(total_row_idx, c)
        cell.fill.solid()
        if c < 4:
            cell.fill.fore_color.rgb = RGBColor(255, 255, 255) # white background
            text_color = logo_color
        else:
            cell.fill.fore_color.rgb = RGBColor(30, 45, 65) # keep original dark bg for grand total
            text_color = TEXT_PRIMARY
            
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.text = text
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = text_color
        p.font.name = 'Arial'
        p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER

    # Add confusion matrix slide
    add_confusion_matrix_slide(prs, inspection_points, notifications)

    # Add issues slides listing details of failed and missing points
    add_issues_slide(prs, failed_points, missing_points)



def find_docx_images(point_name, inspection_type):
    images = {}
    if inspection_type == 'thermal_inspection':
        for key, suffix in [('rgb', '_raw_rgb_image.jpg'), ('thermal', '_thermal_map.jpg')]:
            pattern = os.path.join(MISSION_DIR, f'*{point_name}*{suffix}')
            matches = glob.glob(pattern)
            if matches:
                images[key] = matches[0]
    elif inspection_type in ['leakage_inspection', 'vibration_inspection']:
        for key, suffix in [('soundmap', '_soundmap_photo.jpg'), ('spect', '_spect_photo.jpg')]:
            pattern = os.path.join(MISSION_DIR, f'*{point_name}*{suffix}')
            matches = glob.glob(pattern)
            if matches:
                images[key] = matches[0]
    else:
        for key, suffix in [('processed', '_processed_rgb_image.jpg'), ('preprocessed', '_preprocessed_rgb_image.jpg')]:
            pattern = os.path.join(MISSION_DIR, f'*{point_name}*{suffix}')
            matches = glob.glob(pattern)
            if matches:
                images[key] = matches[0]
    return images

def generate_docx(output_docx_path, inspection_points, notifications, include_all=False):
    doc = Document()
    
    # Page setup (margins: 1.0 inch all around)
    section = doc.sections[0]
    section.page_width = DocxInches(8.5)
    section.page_height = DocxInches(11.0)
    section.top_margin = DocxInches(1.0)
    section.bottom_margin = DocxInches(1.0)
    section.left_margin = DocxInches(1.0)
    section.right_margin = DocxInches(1.0)

    # Helper to set background
    def set_cell_background(cell, fill_hex):
        tcPr = cell._tc.get_or_add_tcPr()
        shd = OxmlElement('w:shd')
        shd.set(qn('w:val'), 'clear')
        shd.set(qn('w:color'), 'auto')
        shd.set(qn('w:fill'), fill_hex)
        tcPr.append(shd)

    # Helper to set cell margins (padding)
    def set_cell_margins(cell, top=144, bottom=144, left=216, right=216):
        tcPr = cell._tc.get_or_add_tcPr()
        tcMar = OxmlElement('w:tcMar')
        for m, val in [('top', top), ('bottom', bottom), ('left', left), ('right', right)]:
            node = OxmlElement(f'w:{m}')
            node.set(qn('w:w'), str(val))
            node.set(qn('w:type'), 'dxa')
            tcMar.append(node)
        tcPr.append(tcMar)

    # Title
    title_p = doc.add_paragraph()
    title_p.paragraph_format.space_before = DocxPt(0)
    title_p.paragraph_format.space_after = DocxPt(4)
    report_title = "Inspection Points Report" if include_all else "Inspection Problem Points Report"
    title_run = title_p.add_run(report_title)
    title_run.font.name = 'Arial'
    title_run.font.size = DocxPt(20)
    title_run.font.bold = True
    title_run.font.color.rgb = DocxRGBColor(31, 41, 55) # #1F2937

    # Subtitle
    sub_p = doc.add_paragraph()
    sub_p.paragraph_format.space_before = DocxPt(0)
    sub_p.paragraph_format.space_after = DocxPt(24)
    sub_run = sub_p.add_run(f"Mission: {MISSION_ID} | Generated automatically from inspection data.")
    sub_run.font.name = 'Arial'
    sub_run.font.size = DocxPt(10)
    sub_run.font.color.rgb = DocxRGBColor(75, 85, 99) # #4B5563

    # Filter or Include all points
    target_points = []
    for point in inspection_points:
        name = point.get('Node_info', '')
        raw_insp = point.get('Inspection', 'General')
        notif = notifications.get(name)
        if notif:
            level = notif.get('notification_level', 'pass').lower()
        else:
            level = 'missing'
            
        if include_all or level != 'pass':
            target_points.append((point, notif or {}, level))
                
    if not target_points:
        msg = "No inspection points detected in this mission." if include_all else "No problem points detected in this mission."
        doc.add_paragraph(msg)
        doc.save(output_docx_path)
        return

    # Add Table
    table = doc.add_table(rows=1, cols=3)
    table.alignment = 1 # Center
    table.allow_autofit = False
    
    # Force fixed table layout in XML to enforce column widths strictly
    tblPr = table._tbl.tblPr
    tblLayout = OxmlElement('w:tblLayout')
    tblLayout.set(qn('w:type'), 'fixed')
    tblPr.append(tblLayout)
    
    # Style Table Borders
    tblBorders = tblPr.first_child_found_in("w:tblBorders")
    if tblBorders is None:
        tblBorders = OxmlElement('w:tblBorders')
        tblPr.append(tblBorders)
    tblBorders.clear()
    for border_name in ['top', 'left', 'bottom', 'right', 'insideH', 'insideV']:
        border = OxmlElement(f'w:{border_name}')
        border.set(qn('w:val'), 'single')
        border.set(qn('w:sz'), '4')
        border.set(qn('w:space'), '0')
        border.set(qn('w:color'), 'E2E8F0')
        tblBorders.append(border)

    # Header Row
    headers = ["Inspection", "Result", "Note/Remark"]
    hdr_cells = table.rows[0].cells
    for c, text in enumerate(headers):
        hdr_cells[c].text = ""
        p = hdr_cells[c].paragraphs[0]
        run = p.add_run(text)
        run.font.name = 'Arial'
        run.font.size = DocxPt(11)
        run.font.bold = True
        run.font.color.rgb = DocxRGBColor(255, 255, 255)
        p.paragraph_format.space_after = DocxPt(4)
        p.paragraph_format.space_before = DocxPt(4)
        set_cell_background(hdr_cells[c], '1F2937')
        set_cell_margins(hdr_cells[c], top=80, bottom=80, left=216, right=216)

    # Column widths: Inspection (~20% = 1.3 in), Result (~50% = 3.25 in), Note/Remark (~30% = 1.95 in)
    col_widths = [DocxInches(1.3), DocxInches(3.25), DocxInches(1.95)]
    for i, col in enumerate(table.columns):
        col.width = col_widths[i]
    for row in table.rows:
        for i, w in enumerate(col_widths):
            row.cells[i].width = w

    # Add data rows
    for idx, (point, notif, level) in enumerate(target_points):
        row_cells = table.add_row().cells
        for i, w in enumerate(col_widths):
            row_cells[i].width = w
            # Column 2 can have slightly smaller left/right padding to prevent image wrapping
            pad_left = 150 if i == 1 else 216
            pad_right = 150 if i == 1 else 216
            set_cell_margins(row_cells[i], top=60, bottom=60, left=pad_left, right=pad_right)

        name = point.get('Node_info', 'N/A')
        raw_insp = point.get('Inspection', 'General')
        type_display = raw_insp.replace('_', ' ').title()
        zone = point.get('Zone', 'N/A')
        floor = point.get('MapName', 'N/A').replace('_', ' ').title()
        csv_res = notif.get('result_dict', {})

        # Set zebra striping backgrounds for columns to match PDF styling
        bg_color = 'FFFFFF' if idx % 2 == 0 else 'F8FAFC'
        set_cell_background(row_cells[0], bg_color)
        set_cell_background(row_cells[1], bg_color)
        set_cell_background(row_cells[2], bg_color)

        # Col 1: Inspection (contains only node name and result metrics)
        p1 = row_cells[0].paragraphs[0]
        p1.paragraph_format.space_before = DocxPt(0)
        p1.paragraph_format.space_after = DocxPt(2)
        
        # Name
        r_name = p1.add_run(f"{name}\n")
        r_name.font.name = 'Arial'
        r_name.font.size = DocxPt(10)
        r_name.font.bold = True
        r_name.font.color.rgb = DocxRGBColor(31, 41, 55) # #1F2937
        
        # Result Metrics List
        res_text = ""
        if raw_insp == 'thermal_inspection':
            threshold = point.get('Threshold', 'N/A')
            min_temp = csv_res.get('min_temperature', 'N/A')
            max_temp = csv_res.get('max_temperature', 'N/A')
            res_text = f"• Status: {level.upper()}\n• Threshold: {threshold}°C\n• Temperature: {min_temp}-{max_temp}°C"
        elif raw_insp == 'leakage_inspection':
            is_leak = csv_res.get('is_leak', 'N/A')
            res_text = f"• Status: {level.upper()}\n• Is Leak: {is_leak}"
        elif raw_insp == 'vibration_inspection':
            k_thresh = point.get('Kurtosis_threshold', 'N/A')
            p_thresh = point.get('Peak_threshold', 'N/A')
            k_val = csv_res.get('kurtosis', 'N/A')
            p_val = csv_res.get('peak_factor', 'N/A')
            res_text = (f"• Status: {level.upper()}\n"
                        f"• Kurtosis Thresh/Value: {k_thresh} / {k_val}\n"
                        f"• Peak Factor Thresh/Value: {p_thresh} / {p_val}")
        else:
            det_status = csv_res.get('detection_status', 'N/A')
            res_text = f"• Status: {level.upper()}\n• Detection Status: {det_status}"
            
        r_res_body = p1.add_run(res_text)
        r_res_body.font.name = 'Arial'
        r_res_body.font.size = DocxPt(8.5)
        r_res_body.font.color.rgb = DocxRGBColor(75, 85, 99) # #4B5563

        # Col 2: Result (contains images side-by-side)
        # Find and insert images
        images = find_docx_images(name, raw_insp)
        img_paths = [path for path in images.values() if os.path.exists(path)]
        
        if len(img_paths) == 1:
            p2 = row_cells[1].paragraphs[0]
            p2.text = ""
            p2.paragraph_format.space_before = DocxPt(0)
            p2.paragraph_format.space_after = DocxPt(0)
            p2.alignment = 1 # Center
            r_img = p2.add_run()
            try:
                r_img.add_picture(img_paths[0], width=DocxInches(1.3), height=DocxInches(0.9))
            except Exception as e:
                print(f"Error adding picture {img_paths[0]}: {e}")
        elif len(img_paths) >= 2:
            # Create a nested table inside the cell
            nested_table = row_cells[1].add_table(rows=1, cols=2)
            nested_table.alignment = 1 # Center
            nested_table.allow_autofit = False
            
            # Remove borders from nested table
            nested_tblPr = nested_table._tbl.tblPr
            nested_tblBorders = OxmlElement('w:tblBorders')
            for border_name in ['top', 'left', 'bottom', 'right', 'insideH', 'insideV']:
                border = OxmlElement(f'w:{border_name}')
                border.set(qn('w:val'), 'none')
                nested_tblBorders.append(border)
            nested_tblPr.append(nested_tblBorders)
            
            nested_widths = [DocxInches(1.48), DocxInches(1.48)]
            for col_idx, col in enumerate(nested_table.columns):
                col.width = nested_widths[col_idx]
                
            n_cells = nested_table.rows[0].cells
            for col_idx in range(2):
                n_cells[col_idx].width = nested_widths[col_idx]
                set_cell_margins(n_cells[col_idx], top=0, bottom=0, left=36, right=36) # very tiny padding
                set_cell_background(n_cells[col_idx], bg_color)
                
                # Add picture
                np_para = n_cells[col_idx].paragraphs[0]
                np_para.text = ""
                np_para.paragraph_format.space_before = DocxPt(0)
                np_para.paragraph_format.space_after = DocxPt(0)
                np_para.alignment = 1 # Center
                
                run = np_para.add_run()
                try:
                    run.add_picture(img_paths[col_idx], width=DocxInches(1.3), height=DocxInches(0.9))
                except Exception as e:
                    print(f"Error adding picture {img_paths[col_idx]}: {e}")
            
            # Delete the default empty paragraph to prevent extra vertical space
            p_default = row_cells[1].paragraphs[0]
            p_default._element.getparent().remove(p_default._element)
        else:
            p2 = row_cells[1].paragraphs[0]
            p2.text = ""
            p2.paragraph_format.space_before = DocxPt(0)
            p2.paragraph_format.space_after = DocxPt(0)

        # Col 3: Note/Remark
        p3 = row_cells[2].paragraphs[0]
        p3.text = ""
        p3.paragraph_format.space_before = DocxPt(0)
        p3.paragraph_format.space_after = DocxPt(0)
        p3.alignment = 3 # Justified alignment
        
        discussion_text = get_discussion(name, raw_insp, csv_res, level.title())
        r_note = p3.add_run(discussion_text)
        r_note.font.name = 'Arial'
        r_note.font.size = DocxPt(9.0)
        r_note.font.color.rgb = DocxRGBColor(75, 85, 99) # #4B5563

    doc.save(output_docx_path)


def main():
    global MISSION_ID, MISSION_DIR, LOGO_PATH
    
    # --- Configuration ---
    MISSION_ID = 'mission-65'
    # json_path = os.path.join(PROJECT_ROOT, 'resource', 'path', 'final_dry_first_floor.json')
    # json_path = os.path.join(PROJECT_ROOT, 'resource', 'path', 'final_dry_first_floor_test.json')
    # json_path = os.path.join(PROJECT_ROOT, 'resource', 'path', 'final_dry_second_floor.json')
    # json_path = os.path.join(PROJECT_ROOT, 'resource', 'path', 'final_packing.json')
    # json_path = os.path.join(PROJECT_ROOT, 'resource', 'path', 'final_packing_test.json')
    # json_path = os.path.join(PROJECT_ROOT, 'resource', 'path', 'final_packing_test_1month.json')
    json_path = os.path.join(PROJECT_ROOT, 'resource', 'path', 'final_filling.json')
    # json_path = os.path.join(PROJECT_ROOT, 'resource', 'path', 'final_filling_1month.json')
    # json_path = os.path.join(PROJECT_ROOT, 'resource', 'path', 'final_filling_1month_old.json')
    
    # Derived Paths
    MISSION_DIR = os.path.join(PROJECT_ROOT, 'resource', MISSION_ID)
    LOGO_PATH = os.path.join(PROJECT_ROOT, 'resource', 'gensurv-logo.jpg')
    output_path = os.path.join(PROJECT_ROOT, f'inspection-{MISSION_ID}.pptx')
    
    load_notifications()
    # ---------------------
    
    if not os.path.exists(json_path):
        print(f"Error: {json_path} not found.")
        return

    with open(json_path, 'r') as f:
        data = json.load(f)
        
    inspection_points = [p for p in data if p.get('PointInfo') == 1 and p.get('Inspection') != 'sit']
    
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH)
    prs.slide_height = Inches(SLIDE_HEIGHT)

    print(f"Found {len(inspection_points)} inspection points. Generating slides...")
    for i, point in enumerate(inspection_points):
        add_slide(prs, point, i, len(inspection_points))
        if (i + 1) % 10 == 0:
            print(f"Generated {i + 1} slides...")
    
    # Add final summary slide
    add_summary_slide(prs, inspection_points, notifications)
    
    prs.save(output_path)
    print(f"Successfully generated {output_path}")

    # Generate Word Document (.docx) report (Original)
    docx_output_path = os.path.join(PROJECT_ROOT, f'inspection-{MISSION_ID}.docx')
    print("Generating Word Document (.docx) report...")
    generate_docx(docx_output_path, inspection_points, notifications, include_all=False)
    print(f"Successfully generated {docx_output_path}")

    # Generate Word Document (.docx) report (All points)
    docx_all_output_path = os.path.join(PROJECT_ROOT, f'inspection-{MISSION_ID}-all.docx')
    print("Generating Word Document (.docx) report for all points...")
    generate_docx(docx_all_output_path, inspection_points, notifications, include_all=True)
    print(f"Successfully generated {docx_all_output_path}")

if __name__ == "__main__":
    main()
